#!/usr/bin/env python3
"""Standalone document parser for stella-meet.
Usage:
  python3 stella_parse.py --file <path> --type pdf|docx|pptx
  python3 stella_parse.py --url <url> --type web
Output: JSON to stdout: {"text": "extracted text content"}
"""

import argparse
import json
import sys


def parse_pdf(path: str) -> str:
    import fitz  # pymupdf
    import os

    # Redirect stdout to suppress library diagnostic messages that corrupt JSON output.
    real_stdout = os.dup(1)
    devnull = os.open(os.devnull, os.O_WRONLY)
    os.dup2(devnull, 1)

    try:
        return _parse_pdf_inner(path, fitz)
    finally:
        # Restore stdout.
        os.dup2(real_stdout, 1)
        os.close(real_stdout)
        os.close(devnull)


def _parse_pdf_inner(path: str, fitz) -> str:
    # First try pymupdf4llm for text-based PDFs.
    try:
        import pymupdf4llm
        text = pymupdf4llm.to_markdown(path)
        if text.strip():
            return text
    except Exception:
        pass

    # Fallback: plain text extraction.
    doc = fitz.open(path)
    text = ""
    for page in doc:
        text += page.get_text() + "\n"
    if text.strip():
        return text

    # Image-only PDF — OCR each page via pytesseract.
    try:
        import io
        import pytesseract
        from PIL import Image

        parts = []
        for page in doc:
            pix = page.get_pixmap(dpi=200)
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            page_text = pytesseract.image_to_string(img)
            if page_text.strip():
                parts.append(page_text)
        if parts:
            return "\n\n".join(parts)
    except ImportError:
        pass

    return ""


def parse_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        if para.text.strip():
            lines.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append(" | ".join(cells))
    return "\n".join(lines)


def parse_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
        lines.append("")
    return "\n".join(lines)


def parse_web(url: str) -> str:
    import trafilatura
    downloaded = trafilatura.fetch_url(url)
    if downloaded is None:
        raise ValueError(f"Failed to fetch URL: {url}")
    text = trafilatura.extract(downloaded)
    if text is None:
        raise ValueError(f"No content extracted from: {url}")
    return text


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--file", help="Path to file")
    parser.add_argument("--url", help="URL to fetch (for web type)")
    parser.add_argument("--type", required=True, choices=["pdf", "docx", "pptx", "web"])
    args = parser.parse_args()

    if args.type == "web":
        if not args.url:
            print(json.dumps({"error": "--url is required for web type"}))
            sys.exit(1)
        try:
            text = parse_web(args.url)
            json.dump({"text": text}, sys.stdout)
        except Exception as e:
            json.dump({"error": str(e)}, sys.stdout)
            sys.exit(1)
    else:
        if not args.file:
            print(json.dumps({"error": "--file is required for file types"}))
            sys.exit(1)
        parsers = {
            "pdf": parse_pdf,
            "docx": parse_docx,
            "pptx": parse_pptx,
        }
        try:
            text = parsers[args.type](args.file)
            json.dump({"text": text}, sys.stdout)
        except Exception as e:
            json.dump({"error": str(e)}, sys.stdout)
            sys.exit(1)


if __name__ == "__main__":
    main()
